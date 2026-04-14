"""
generate_ppt.py
Creates a professional leadership presentation for GitAnalytics AI Dashboard.
Run: python3 generate_ppt.py
Output: GitAnalytics_AI_Dashboard_Leadership.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ── Brand Palette ─────────────────────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
CARD_BG       = RGBColor(0x1E, 0x29, 0x3B)   # Dark card
ACCENT_TEAL   = RGBColor(0x2D, 0xD4, 0xBF)   # Manual LOC
ACCENT_VIOLET = RGBColor(0xA7, 0x8B, 0xFA)   # AI LOC
ACCENT_BLUE   = RGBColor(0x60, 0xA5, 0xFA)   # Brand blue
ACCENT_GREEN  = RGBColor(0x34, 0xD3, 0x99)   # Success/positive
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)   # Warning
ACCENT_RED    = RGBColor(0xF8, 0x71, 0x71)   # Danger/AI flag
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0x94, 0xA3, 0xB8)
MID_GREY      = RGBColor(0x4B, 0x55, 0x63)
DARKER_CARD   = RGBColor(0x11, 0x1B, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]   # Completely blank

# ─── Helper Functions ─────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def fill_bg(slide, color=DARK_BG):
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def box(slide, l, t, w, h, color=CARD_BG, radius=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def accent_bar(slide, l, t, w, h=0.05, color=ACCENT_BLUE):
    return box(slide, l, t, w, h, color)

def text_box(slide, text, l, t, w, h,
             font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def multi_para(slide, lines, l, t, w, h,
               font_size=12, color=WHITE, bold_first=False,
               spacing_after=6, leading_color=None, line_spacing=1.1):
    """Add multiple lines as separate paragraphs in one text box."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        use_color = leading_color if (leading_color and i == 0) else color
        run.font.color.rgb = use_color
    return txBox

def pill(slide, label, l, t, w=1.5, h=0.35, bg=ACCENT_TEAL, fg=DARK_BG, fs=10):
    """Rounded pill / badge."""
    shape = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(fs)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape

def icon_card(slide, icon, title, desc, l, t, w=2.8, h=1.5,
              icon_color=ACCENT_TEAL, title_color=WHITE):
    """Small icon card with title and description."""
    card = box(slide, l, t, w, h, CARD_BG)
    # top colour accent strip
    accent_bar(slide, l, t, w, 0.05, icon_color)
    # icon text
    text_box(slide, icon, l+0.18, t+0.15, 0.5, 0.45,
             font_size=20, bold=True, color=icon_color)
    # title
    text_box(slide, title, l+0.65, t+0.18, w-0.75, 0.35,
             font_size=12, bold=True, color=title_color)
    # description
    text_box(slide, desc, l+0.18, t+0.65, w-0.28, 0.75,
             font_size=10, color=LIGHT_GREY)
    return card

def kpi_card(slide, value, label, l, t, w=2.1, h=1.3,
             value_color=ACCENT_TEAL, bg=CARD_BG):
    card = box(slide, l, t, w, h, bg)
    accent_bar(slide, l, t, w, 0.05, value_color)
    text_box(slide, value, l+0.15, t+0.2, w-0.2, 0.6,
             font_size=28, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    text_box(slide, label, l+0.1, t+0.82, w-0.2, 0.42,
             font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    return card

def slide_header(slide, title, subtitle=None, accent_color=ACCENT_BLUE):
    accent_bar(slide, 0, 0, 13.33, 0.07, accent_color)
    text_box(slide, title, 0.5, 0.15, 12.0, 0.55,
             font_size=22, bold=True, color=WHITE)
    if subtitle:
        text_box(slide, subtitle, 0.5, 0.68, 12.0, 0.38,
                 font_size=13, color=LIGHT_GREY, italic=True)
    # horizontal rule
    rule = slide.shapes.add_shape(1,
        Inches(0.5), Inches(1.12), Inches(12.33), Inches(0.018))
    rule.fill.solid()
    rule.fill.fore_color.rgb = MID_GREY
    rule.line.fill.background()

def flow_arrow(slide, x, y, w=0.5, vertical=False, color=ACCENT_BLUE):
    """A simple rectangle arrow indicator."""
    if vertical:
        box(slide, x, y, 0.04, w, color)
        # arrowhead triangle
        tri = slide.shapes.add_shape(7,
            Inches(x-0.06), Inches(y+w-0.02), Inches(0.16), Inches(0.18))
        tri.fill.solid(); tri.fill.fore_color.rgb = color
        tri.line.fill.background()
    else:
        box(slide, x, y, w, 0.04, color)
        tri = slide.shapes.add_shape(8,
            Inches(x+w-0.02), Inches(y-0.06), Inches(0.18), Inches(0.16))
        tri.fill.solid(); tri.fill.fore_color.rgb = color
        tri.line.fill.background()

def progress_bar(slide, l, t, w, h, pct, bg=MID_GREY,
                 fg=ACCENT_TEAL, label=None):
    box(slide, l, t, w, h, bg)
    box(slide, l, t, w*pct, h, fg)
    if label:
        text_box(slide, label, l+w+0.1, t-0.02, 0.6, h+0.05,
                 font_size=10, bold=True, color=fg)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide, DARKER_CARD)

# Large geometric accent shape (top-right)
corner = slide.shapes.add_shape(1,
    Inches(9.5), Inches(-0.5), Inches(5.0), Inches(5.0))
corner.fill.solid()
corner.fill.fore_color.rgb = RGBColor(0x1A, 0x26, 0x3D)
corner.line.fill.background()

# Teal glow strip on the left
accent_bar(slide, 0, 0, 0.12, 7.5, ACCENT_TEAL)

# Decorative dots grid
for row in range(5):
    for col in range(6):
        dot = slide.shapes.add_shape(9,
            Inches(9.8+col*0.5), Inches(0.5+row*0.5),
            Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(0x2D, 0x3B, 0x55)
        dot.line.fill.background()

# Product badge
pill(slide, "  v1.0  LEADERSHIP BRIEF", 0.5, 1.0, 3.0, 0.38,
     ACCENT_BLUE, WHITE, 10)

# Main title
text_box(slide,
    "GitAnalytics AI\nDashboard",
    0.5, 1.6, 9.0, 1.8,
    font_size=52, bold=True, color=WHITE)

# Subtitle
text_box(slide,
    "Quantifying AI-Assisted Code Contributions\nAcross Your Engineering Teams",
    0.5, 3.55, 9.0, 1.0,
    font_size=20, color=ACCENT_TEAL, italic=True)

# Description paragraph
text_box(slide,
    "A read-only, browser-based analytics platform that connects to GitLab to\n"
    "automatically detect, measure, and visualise AI vs. manually written code —\n"
    "giving leadership real-time insight into AI adoption across every project.",
    0.5, 4.7, 8.5, 1.2,
    font_size=13, color=LIGHT_GREY)

# Tag pills
tags = ["React 18", "TypeScript", "Recharts", "GitLab REST API", "Zero Backend"]
for i, tag in enumerate(tags):
    pill(slide, tag, 0.5+i*2.1, 6.1, 2.0, 0.38,
         CARD_BG, ACCENT_BLUE, 10)

# Date
text_box(slide, "April 2026", 0.5, 6.95, 3.0, 0.4,
         font_size=11, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Executive Summary",
             "What it does · Why it matters · How it works", ACCENT_TEAL)

# Left column - problem/solution
box(slide, 0.4, 1.3, 6.0, 5.6, CARD_BG)
accent_bar(slide, 0.4, 1.3, 6.0, 0.06, ACCENT_ORANGE)
text_box(slide, "The Challenge", 0.65, 1.42, 5.5, 0.45,
         font_size=15, bold=True, color=ACCENT_ORANGE)
multi_para(slide, [
    "Engineering leaders cannot easily answer:",
    "",
    "  •  How much of our code is AI-generated?",
    "  •  Which developers are leveraging AI tools?",
    "  •  Is AI adoption improving team velocity?",
    "  •  Are AI contribution patterns healthy?",
    "",
    "Manual audits are time-consuming, inconsistent,",
    "and require access to every commit across",
    "multiple GitLab projects — infeasible at scale.",
], 0.65, 1.92, 5.5, 4.8,
font_size=12, color=LIGHT_GREY)

# Right column - solution points
box(slide, 6.9, 1.3, 6.0, 5.6, CARD_BG)
accent_bar(slide, 6.9, 1.3, 6.0, 0.06, ACCENT_GREEN)
text_box(slide, "The Solution", 7.15, 1.42, 5.5, 0.45,
         font_size=15, bold=True, color=ACCENT_GREEN)

solution_items = [
    ("🔍", "Automated AI Detection",
     "3 independent signals detect AI-generated commits without manual tagging."),
    ("📊", "Real-Time LOC Metrics",
     "Lines of code split by Manual vs AI per developer, per project."),
    ("🛡️", "Zero Data Risk",
     "No backend server. Token never persisted. Data cleared on page refresh."),
    ("⚡", "Instant Deployment",
     "Static SPA — runs from any CDN, S3 bucket, or web server. No DevOps needed."),
]
for i, (icon, title, desc) in enumerate(solution_items):
    top = 1.95 + i * 1.22
    box(slide, 7.1, top, 5.6, 1.05, RGBColor(0x14, 0x20, 0x35))
    text_box(slide, icon, 7.25, top+0.15, 0.5, 0.6, font_size=18)
    text_box(slide, title, 7.8, top+0.1, 4.7, 0.38,
             font_size=12, bold=True, color=WHITE)
    text_box(slide, desc, 7.8, top+0.48, 4.7, 0.52,
             font_size=10, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM IN NUMBERS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "The AI Adoption Visibility Gap",
             "Why traditional code review tools miss the picture", ACCENT_VIOLET)

# Large stat cards
stats = [
    ("77%", "of developers", "now use AI coding tools\nat least weekly",
     ACCENT_VIOLET),
    ("3×", "faster feature", "delivery reported with\nAI-assisted workflows",
     ACCENT_TEAL),
    ("0%", "of git dashboards", "natively classify\nAI vs manual LOC",
     ACCENT_ORANGE),
    ("100%", "manual today", "leadership relies on\nself-reported surveys",
     ACCENT_RED),
]
for i, (val, sub1, sub2, color) in enumerate(stats):
    l = 0.5 + i * 3.2
    box(slide, l, 1.35, 3.0, 2.5, CARD_BG)
    accent_bar(slide, l, 1.35, 3.0, 0.06, color)
    text_box(slide, val, l, 1.6, 3.0, 0.9,
             font_size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
    text_box(slide, sub1, l, 2.52, 3.0, 0.38,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(slide, sub2, l, 2.9, 3.0, 0.85,
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Gap illustration
box(slide, 0.5, 4.1, 12.33, 2.8, CARD_BG)
text_box(slide, "Without GitAnalytics", 0.8, 4.25, 5.7, 0.45,
         font_size=14, bold=True, color=ACCENT_RED)
text_box(slide, "With GitAnalytics", 7.0, 4.25, 5.7, 0.45,
         font_size=14, bold=True, color=ACCENT_GREEN)

# VS divider
text_box(slide, "VS", 6.3, 4.6, 0.8, 0.8,
         font_size=22, bold=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
rule = slide.shapes.add_shape(1,
    Inches(6.6), Inches(4.15), Inches(0.04), Inches(2.6))
rule.fill.solid(); rule.fill.fore_color.rgb = MID_GREY
rule.line.fill.background()

without = [
    "✗  No visibility into AI adoption rate",
    "✗  Surveys are slow & self-reporting biased",
    "✗  No per-developer breakdown",
    "✗  Impossible to trend over time",
]
with_ = [
    "✓  Instant AI/manual LOC split per dev",
    "✓  3 automated detection signals",
    "✓  Drill-down to individual commits",
    "✓  Compare periods with date range filter",
]
multi_para(slide, without, 0.8, 4.75, 5.5, 1.9,
           font_size=11, color=LIGHT_GREY)
multi_para(slide, with_, 7.0, 4.75, 5.5, 1.9,
           font_size=11, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Technology Stack",
             "Modern, lightweight, and zero-dependency on external services",
             ACCENT_BLUE)

layers = [
    ("PRESENTATION LAYER", [
        ("React 18.2", "Component library\n(Hooks, Strict Mode)", ACCENT_BLUE),
        ("TypeScript 5.3", "Static typing\n(strict mode)", ACCENT_VIOLET),
        ("Tailwind CSS 3.4", "Utility-first styling\n+ brand theme", ACCENT_TEAL),
    ]),
    ("VISUALISATION LAYER", [
        ("Recharts 2.10", "Composable React charts\n(Bar, Pie, Tooltip)", ACCENT_ORANGE),
        ("Custom Tooltips", "Hover interactions\nwith K-suffix formatting", ACCENT_BLUE),
        ("Donut + Bar Views", "Toggle between chart\nperspectives", ACCENT_TEAL),
    ]),
    ("BUILD & DELIVERY", [
        ("Vite 5.0", "Lightning-fast bundler\nHMR dev server :3000", ACCENT_GREEN),
        ("PostCSS / Autoprefixer", "CSS transformation\npipeline", ACCENT_VIOLET),
        ("Static SPA Output", "dist/ → any CDN\nor web server", ACCENT_TEAL),
    ]),
]

for col_idx, (layer_title, cards) in enumerate(layers):
    lx = 0.45 + col_idx * 4.3
    text_box(slide, layer_title, lx, 1.35, 4.0, 0.35,
             font_size=9, bold=True, color=ACCENT_BLUE,
             align=PP_ALIGN.CENTER)
    # layer bracket
    box(slide, lx, 1.68, 4.0, 0.03, MID_GREY)
    for row_idx, (name, desc, color) in enumerate(cards):
        ty = 1.8 + row_idx * 1.72
        box(slide, lx, ty, 4.0, 1.55, CARD_BG)
        accent_bar(slide, lx, ty, 4.0, 0.05, color)
        text_box(slide, name, lx+0.2, ty+0.12, 3.6, 0.4,
                 font_size=14, bold=True, color=color)
        text_box(slide, desc, lx+0.2, ty+0.55, 3.6, 0.9,
                 font_size=11, color=LIGHT_GREY)

# Bottom note
text_box(slide,
    "No backend server required · No database · No authentication service · Zero server-side code",
    0.5, 7.1, 12.33, 0.35,
    font_size=11, color=ACCENT_GREEN, align=PP_ALIGN.CENTER, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI DETECTION ENGINE (THE CORE INNOVATION)
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "AI Detection Engine",
             "Three independent signals — any match classifies the commit as AI-assisted",
             ACCENT_VIOLET)

# Centre pipeline
signals = [
    ("01", "Commit Message\nPattern Matching",
     "18 regex patterns scanning\ntitle + body for AI tool\nconventions",
     ACCENT_VIOLET,
     ["[AI-Generated]", "[Claude]", "copilot:", "[ChatGPT]",
      "[Cursor]", "Co-Authored-By: *bot*"]),
    ("02", "Committer Email\nMismatch",
     "Flags when committer ≠ author\n— tool committed on\nbehalf of developer",
     ACCENT_TEAL,
     ["GitHub Copilot auto-commit", "Cursor autonomous mode",
      "CI bot commits", "Agent-driven push"]),
    ("03", "LOC Burst\nHeuristic",
     "Flags manual commits\n≥ 3× developer's own\nmedian additions",
     ACCENT_ORANGE,
     ["Requires ≥ 5 commits", "Developer-specific baseline",
      "Catches undeclared AI use", "Statistical anomaly detection"]),
]

for i, (num, title, desc, color, examples) in enumerate(signals):
    lx = 0.4 + i * 4.3
    # Card
    card = box(slide, lx, 1.35, 4.0, 5.6, CARD_BG)
    accent_bar(slide, lx, 1.35, 4.0, 0.07, color)

    # Number badge
    pill(slide, num, lx+0.2, 1.5, 0.55, 0.38, color, DARK_BG, 12)
    # Title
    text_box(slide, title, lx+0.9, 1.48, 3.0, 0.5,
             font_size=13, bold=True, color=color)
    # Description
    text_box(slide, desc, lx+0.2, 2.08, 3.6, 0.9,
             font_size=11, color=LIGHT_GREY)

    # Examples box
    box(slide, lx+0.2, 3.1, 3.6, 2.6, RGBColor(0x14, 0x20, 0x35))
    text_box(slide, "EXAMPLES", lx+0.35, 3.17, 3.2, 0.3,
             font_size=8, bold=True, color=MID_GREY)
    for j, ex in enumerate(examples):
        text_box(slide, f"  › {ex}", lx+0.28, 3.5+j*0.44, 3.4, 0.42,
                 font_size=10, color=LIGHT_GREY)

# Result badge at bottom
box(slide, 3.0, 7.0, 7.33, 0.42, RGBColor(0x1A, 0x3A, 0x2E))
text_box(slide,
    "  Any signal match  →  isAiGenerated = true  |  aiReason tagged for audit trail",
    3.0, 7.0, 7.33, 0.42,
    font_size=11, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE & DATA FLOW
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Architecture & Data Flow",
             "Client-side SPA — all processing in the browser, GitLab is the only external dependency",
             ACCENT_BLUE)

# 5 pipeline stages
stages = [
    ("🔧", "Configure",
     "User enters\nGitLab URL,\nPAT token,\nGroup/Project IDs",
     ACCENT_BLUE),
    ("📡", "Fetch",
     "GitLabService\ncalls REST v4\nwith pagination\n& rate-limit retry",
     ACCENT_TEAL),
    ("⚙️", "Process",
     "locParser runs\n3-signal AI\ndetection &\nLOC aggregation",
     ACCENT_VIOLET),
    ("📊", "Visualise",
     "Recharts renders\nKPI cards, bar\nchart, pie chart\nper developer",
     ACCENT_ORANGE),
    ("🔎", "Drill-Down",
     "Click any dev row\nto inspect individual\ncommits with\nAI reason badges",
     ACCENT_GREEN),
]

stage_y = 1.5
for i, (icon, name, desc, color) in enumerate(stages):
    lx = 0.3 + i * 2.58
    # stage box
    box(slide, lx, stage_y, 2.3, 3.6, CARD_BG)
    accent_bar(slide, lx, stage_y, 2.3, 0.06, color)
    text_box(slide, icon, lx+0.9, stage_y+0.12, 0.5, 0.5, font_size=22)
    text_box(slide, name, lx+0.05, stage_y+0.65, 2.2, 0.4,
             font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    text_box(slide, desc, lx+0.1, stage_y+1.12, 2.1, 2.2,
             font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    # connector arrow (not after last)
    if i < 4:
        arrow_x = lx + 2.35
        box(slide, arrow_x, stage_y+1.6, 0.2, 0.04, ACCENT_BLUE)
        # arrowhead
        tri = slide.shapes.add_shape(8,
            Inches(arrow_x+0.17), Inches(stage_y+1.53),
            Inches(0.14), Inches(0.18))
        tri.fill.solid(); tri.fill.fore_color.rgb = ACCENT_BLUE
        tri.line.fill.background()

# Bottom detail row
details = [
    ("ConfigPanel.tsx", "collapsible form, token hidden by default", ACCENT_BLUE),
    ("GitLabService.ts", "REST v4, pagination 100/page, exp back-off", ACCENT_TEAL),
    ("locParser.ts", "aggregateByDeveloper() — 2-pass algorithm", ACCENT_VIOLET),
    ("CommitChart.tsx", "Bar (stacked) + Pie (donut) with toggle", ACCENT_ORANGE),
    ("StatsTable.tsx", "expandable rows, 10 commits/page pager", ACCENT_GREEN),
]
for i, (file, note, color) in enumerate(details):
    lx = 0.3 + i * 2.58
    text_box(slide, file, lx, 5.35, 2.3, 0.32,
             font_size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
    text_box(slide, note, lx, 5.65, 2.3, 0.55,
             font_size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Token flow note
box(slide, 0.3, 6.3, 12.73, 0.95, RGBColor(0x14, 0x20, 0x35))
text_box(slide, "SECURITY NOTE", 0.6, 6.38, 2.0, 0.3,
         font_size=9, bold=True, color=ACCENT_ORANGE)
text_box(slide,
    "Token flows: UI → React state → GitLabService (function arg) → GitLab HTTPS request\n"
    "Token NEVER written to localStorage / cookies / server logs. Cleared on page refresh.",
    2.7, 6.35, 10.1, 0.85,
    font_size=10, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — KEY FEATURES & UI WALKTHROUGH
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Dashboard Features",
             "Every view designed for engineering leadership decision-making",
             ACCENT_TEAL)

features = [
    ("📈", "KPI Cards",
     "5 live metrics: Active Developers, Total Commits, "
     "Total LOC, Manual LOC, AI LOC — updated on every fetch.",
     ACCENT_TEAL),
    ("📉", "LOC Bar Chart",
     "Per-developer stacked bar chart "
     "(Manual teal / AI violet). Y-axis auto-scales with K suffix for large values.",
     ACCENT_VIOLET),
    ("🍩", "Global Donut",
     "Overall AI% donut with legend showing raw line counts. "
     "One-click toggle between Bar and Pie views.",
     ACCENT_ORANGE),
    ("🗂️", "Developer Table",
     "Sortable breakdown: name, commit counts, "
     "LOC columns, and a mini stacked progress bar for AI split at a glance.",
     ACCENT_BLUE),
    ("🔍", "Commit Drill-Down",
     "Expand any developer row to see all commits: "
     "SHA, title, date, additions, deletions, and AI reason badge.",
     ACCENT_GREEN),
    ("🏷️", "Author Filter",
     "Tag-style multi-author filter. "
     "Case-insensitive substring match on name or email. Backspace to remove.",
     ACCENT_ORANGE),
    ("📅", "Date Range Picker",
     "Quick presets (7d / 14d / 30d / 90d) "
     "plus manual From/To date inputs for custom sprint or quarter views.",
     ACCENT_TEAL),
    ("⚙️", "Whitespace Discount",
     "Toggle a 5% LOC reduction that removes "
     "blank-line and formatting-only inflation from raw addition counts.",
     ACCENT_VIOLET),
]

cols, rows = 4, 2
card_w, card_h = 3.0, 2.5
for idx, (icon, title, desc, color) in enumerate(features):
    col = idx % cols
    row = idx // cols
    lx = 0.38 + col * 3.22
    ty = 1.35 + row * 2.7
    card = box(slide, lx, ty, card_w, card_h, CARD_BG)
    accent_bar(slide, lx, ty, card_w, 0.05, color)
    text_box(slide, icon, lx+0.15, ty+0.12, 0.5, 0.45, font_size=18)
    text_box(slide, title, lx+0.65, ty+0.15, card_w-0.75, 0.38,
             font_size=12, bold=True, color=color)
    text_box(slide, desc, lx+0.15, ty+0.6, card_w-0.25, 1.75,
             font_size=10, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LIVE KPI MOCKUP
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Sample Dashboard Output",
             "Illustrative example — real data comes live from your GitLab instance",
             ACCENT_BLUE)

# Simulate KPI bar
kpis = [
    ("12", "Active Developers", ACCENT_TEAL),
    ("847", "Total Commits", ACCENT_BLUE),
    ("142,380", "Total LOC", ACCENT_VIOLET),
    ("88,550", "Manual LOC", ACCENT_GREEN),
    ("53,830", "AI-Assisted LOC", ACCENT_ORANGE),
]
for i, (val, label, color) in enumerate(kpis):
    kpi_card(slide, val, label, 0.38+i*2.58, 1.35, 2.4, 1.3, color)

# AI% gauge label
text_box(slide, "37.8% AI-Assisted", 4.5, 2.75, 4.33, 0.5,
         font_size=14, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
progress_bar(slide, 0.38, 2.85, 12.55, 0.18, 0.378,
             MID_GREY, ACCENT_ORANGE, "37.8%")

# Simulated bar chart
chart_title = "LOC by Developer (Manual vs AI)"
text_box(slide, chart_title, 0.38, 3.25, 7.5, 0.38,
         font_size=13, bold=True, color=WHITE)
box(slide, 0.38, 3.65, 7.5, 3.3, CARD_BG)

devs = [
    ("alice@co.com",    5800, 2900),
    ("bob@co.com",      4200, 3100),
    ("carol@co.com",    6100, 800),
    ("dave@co.com",     3300, 3300),
    ("elena@co.com",    2700, 400),
]
max_loc = max(m+a for _,m,a in devs)
bar_area_w = 6.5
bar_area_h = 2.6
bar_area_l = 0.7
bar_area_t = 3.8
bar_w = bar_area_w / len(devs) * 0.55
gap   = bar_area_w / len(devs)

for i, (name, manual, ai) in enumerate(devs):
    bx = bar_area_l + i*gap + gap*0.2
    total = manual + ai
    total_h = (total / max_loc) * bar_area_h
    manual_h = (manual / max_loc) * bar_area_h
    ai_h = (ai / max_loc) * bar_area_h
    base_y = bar_area_t + bar_area_h - total_h
    box(slide, bx, base_y, bar_w, manual_h, ACCENT_TEAL)
    box(slide, bx, base_y+manual_h, bar_w, ai_h, ACCENT_VIOLET)
    short = name.split("@")[0]
    text_box(slide, short, bx-0.1, bar_area_t+bar_area_h+0.05, bar_w+0.2, 0.3,
             font_size=9, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Legend
pill(slide, "  Manual LOC", 0.7, 6.7, 1.9, 0.32, ACCENT_TEAL, DARK_BG, 9)
pill(slide, "  AI LOC",     2.7, 6.7, 1.6, 0.32, ACCENT_VIOLET, DARK_BG, 9)

# Simulated table
table_l = 8.1
text_box(slide, "Developer Breakdown", table_l, 3.25, 5.1, 0.38,
         font_size=13, bold=True, color=WHITE)
box(slide, table_l, 3.65, 5.1, 3.3, CARD_BG)

headers = ["Developer", "Commits", "Total LOC", "AI %"]
widths  = [1.8, 0.8, 1.1, 1.2]
hx = table_l+0.1
for h, w in zip(headers, widths):
    text_box(slide, h, hx, 3.72, w, 0.3,
             font_size=9, bold=True, color=LIGHT_GREY)
    hx += w

rows_data = [
    ("alice@co.com",   "38",  "8,700",   "33%", 0.33),
    ("bob@co.com",     "51",  "7,300",   "42%", 0.42),
    ("carol@co.com",   "29",  "6,900",   "12%", 0.12),
    ("dave@co.com",    "44",  "6,600",   "50%", 0.50),
    ("elena@co.com",   "22",  "3,100",   "13%", 0.13),
]
for r_idx, (dev, commits, total, ai_pct, ai_raw) in enumerate(rows_data):
    ry = 4.1 + r_idx * 0.55
    if r_idx % 2 == 0:
        box(slide, table_l, ry, 5.1, 0.52, RGBColor(0x14, 0x20, 0x35))
    row_vals = [dev.split("@")[0], commits, total]
    rx = table_l + 0.1
    for v, w in zip(row_vals, widths[:-1]):
        text_box(slide, v, rx, ry+0.1, w, 0.35,
                 font_size=10, color=LIGHT_GREY)
        rx += w
    # mini progress bar
    progress_bar(slide, rx, ry+0.2, 0.9, 0.14, ai_raw,
                 MID_GREY, ACCENT_VIOLET)
    text_box(slide, ai_pct, rx+0.95, ry+0.1, 0.3, 0.3,
             font_size=9, bold=True, color=ACCENT_VIOLET)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — COMMIT DRILL-DOWN DETAIL
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Commit-Level Drill-Down",
             "Click any developer row to inspect individual commits with AI reason tags",
             ACCENT_VIOLET)

# Expanded row mockup
box(slide, 0.4, 1.35, 12.53, 5.95, CARD_BG)

# Dev header row (expanded)
box(slide, 0.4, 1.35, 12.53, 0.65, RGBColor(0x1A, 0x26, 0x3D))
text_box(slide, "▼  bob  ·  bob@company.com", 0.7, 1.45, 4.0, 0.45,
         font_size=13, bold=True, color=ACCENT_TEAL)
kpi_cols = [("51", "Commits"), ("7,300", "Total LOC"),
            ("4,200", "Manual"), ("3,100", "AI")]
for i, (v, l) in enumerate(kpi_cols):
    text_box(slide, v, 5.5+i*1.8, 1.4, 1.5, 0.3,
             font_size=14, bold=True, color=ACCENT_VIOLET, align=PP_ALIGN.CENTER)
    text_box(slide, l, 5.5+i*1.8, 1.68, 1.5, 0.25,
             font_size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Column headers
col_headers = ["Short SHA", "Title", "Date", "+Lines", "−Lines", "Detection"]
col_widths  = [0.95, 4.5, 1.5, 0.9, 0.9, 2.5]
cx = 0.65
for h, w in zip(col_headers, col_widths):
    text_box(slide, h, cx, 2.1, w, 0.3,
             font_size=9, bold=True, color=LIGHT_GREY)
    cx += w

# Commit rows
commits = [
    ("a1b2c3d", "[Claude] Refactor auth middleware to use JWT refresh tokens",
     "Apr 10", "412", "88",  "pattern", ACCENT_VIOLET),
    ("e4f5a6b", "Add unit tests for token validation edge cases",
     "Apr 10", "134", "12",  "manual", ACCENT_GREEN),
    ("7c8d9e0", "co-pilot: optimise database query with index hints",
     "Apr 09", "67",  "5",   "pattern", ACCENT_VIOLET),
    ("f1a2b3c", "Fix null pointer in session handler",
     "Apr 09", "23",  "2",   "manual", ACCENT_GREEN),
    ("d4e5f60", "Implement bulk export endpoint (1,840 lines)",
     "Apr 08", "1840","0",   "loc-burst", ACCENT_ORANGE),
    ("1b2c3d4", "Update CI pipeline for node 20 upgrade",
     "Apr 07", "45",  "38",  "committer-mismatch", ACCENT_TEAL),
    ("5e6f7a8", "Refactor user profile components",
     "Apr 07", "289", "102", "manual", ACCENT_GREEN),
    ("9b0c1d2", "[Cursor] Generate openAPI spec from route decorators",
     "Apr 06", "728", "14",  "pattern", ACCENT_VIOLET),
]

reason_labels = {
    "pattern":            ("PATTERN",  ACCENT_VIOLET),
    "loc-burst":          ("LOC BURST",ACCENT_ORANGE),
    "committer-mismatch": ("TOOL COMMIT", ACCENT_TEAL),
    "manual":             ("MANUAL",   ACCENT_GREEN),
}

for r_idx, (sha, title, date, plus, minus, reason, _color) in enumerate(commits):
    ry = 2.5 + r_idx * 0.54
    if r_idx % 2 == 0:
        box(slide, 0.4, ry, 12.53, 0.52, RGBColor(0x14, 0x20, 0x35))
    row_vals = [sha, title, date, f"+{plus}", f"-{minus}"]
    rx = 0.65
    colors_v = [LIGHT_GREY, WHITE, LIGHT_GREY, ACCENT_GREEN, ACCENT_RED]
    for v, w, cv in zip(row_vals, col_widths[:-1], colors_v):
        text_box(slide, v, rx, ry+0.1, w, 0.35,
                 font_size=10, color=cv)
        rx += w
    r_label, r_color = reason_labels[reason]
    pill(slide, r_label, rx, ry+0.08, col_widths[-1]-0.1, 0.32,
         r_color, DARK_BG if r_color != DARK_BG else WHITE, 9)

# Legend at bottom
legend_items = [
    ("PATTERN",        ACCENT_VIOLET, "AI tool tag in commit message"),
    ("LOC BURST",      ACCENT_ORANGE, "Commit ≥ 3× developer's median size"),
    ("TOOL COMMIT",    ACCENT_TEAL,   "Committer email ≠ author email"),
    ("MANUAL",         ACCENT_GREEN,  "No AI signal detected"),
]
for i, (label, color, desc) in enumerate(legend_items):
    lx = 0.5 + i * 3.2
    pill(slide, label, lx, 7.1, 1.5, 0.32, color, DARK_BG, 9)
    text_box(slide, desc, lx+1.6, 7.12, 1.55, 0.3,
             font_size=9, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SECURITY MODEL
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Security & Privacy Model",
             "Designed for zero data exposure — every byte stays between the browser and GitLab",
             ACCENT_GREEN)

# Token lifecycle diagram
text_box(slide, "Personal Access Token Lifecycle", 0.5, 1.38, 6.0, 0.4,
         font_size=14, bold=True, color=WHITE)

token_stages = [
    ("User Enters\nToken in UI",   ACCENT_BLUE),
    ("Stored in\nReact State\n(RAM only)",  ACCENT_TEAL),
    ("Passed as\nFunction Arg\nto Service", ACCENT_VIOLET),
    ("Sent as Header\nto GitLab\nHTTPS only", ACCENT_ORANGE),
    ("Page Refresh\nClears Token\nfrom Memory", ACCENT_GREEN),
]
for i, (label, color) in enumerate(token_stages):
    lx = 0.4 + i * 2.52
    box(slide, lx, 1.85, 2.2, 1.35, CARD_BG)
    accent_bar(slide, lx, 1.85, 2.2, 0.06, color)
    text_box(slide, label, lx+0.1, 2.0, 2.0, 1.0,
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        box(slide, lx+2.22, 2.42, 0.28, 0.04, color)
        tri = slide.shapes.add_shape(8,
            Inches(lx+2.47), Inches(2.35),
            Inches(0.13), Inches(0.17))
        tri.fill.solid(); tri.fill.fore_color.rgb = color
        tri.line.fill.background()

# What IS and IS NOT transmitted
col_configs = [
    ("NEVER STORED OR TRANSMITTED", ACCENT_RED, [
        "✗  localStorage / sessionStorage",
        "✗  cookies or IndexedDB",
        "✗  console.log or browser devtools",
        "✗  any analytics / telemetry service",
        "✗  any backend server",
        "✗  the git repository or CI/CD logs",
    ]),
    ("ONLY TRANSMITTED TO GITLAB", ACCENT_GREEN, [
        "✓  PRIVATE-TOKEN header in HTTPS requests",
        "✓  Only to the base URL the user provided",
        "✓  GitLab REST v4 endpoints only",
        "✓  Cleared from memory after each request",
        "✓  No caching at network or app layer",
        "✓  Read-only scopes only (read_api)",
    ]),
    ("DATA RETURNED FROM GITLAB", ACCENT_BLUE, [
        "•  Commit SHAs, titles, messages",
        "•  Author & committer names/emails",
        "•  Additions & deletions counts",
        "•  All held in React state (RAM)",
        "•  Cleared on page refresh",
        "•  No server-side persistence",
    ]),
]
for i, (title, color, items) in enumerate(col_configs):
    lx = 0.4 + i * 4.3
    box(slide, lx, 3.45, 4.0, 3.75, CARD_BG)
    accent_bar(slide, lx, 3.45, 4.0, 0.06, color)
    text_box(slide, title, lx+0.15, 3.55, 3.7, 0.38,
             font_size=10, bold=True, color=color)
    multi_para(slide, items, lx+0.15, 4.0, 3.7, 3.1,
               font_size=11, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — LOC METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "LOC Measurement Methodology",
             "How we count, what we include, and the tradeoffs made",
             ACCENT_ORANGE)

# Decision tree diagram
box(slide, 0.4, 1.35, 12.53, 5.9, CARD_BG)

text_box(slide, "RAW COMMIT", 5.5, 1.55, 2.33, 0.45,
         font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(slide, 5.5, 1.55, 2.33, 0.4, RGBColor(0x1A, 0x26, 0x3D))

# Merge check
box(slide, 4.8, 2.25, 3.73, 0.55, MID_GREY)
text_box(slide, "Is merge commit? (parent_ids > 1)", 4.9, 2.3, 3.55, 0.4,
         font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Yes branch - excluded
box(slide, 1.0, 3.1, 2.5, 0.55, ACCENT_RED)
text_box(slide, "EXCLUDED\n(prevents double-count)", 1.0, 3.1, 2.5, 0.55,
         font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text_box(slide, "YES →", 1.8, 2.7, 2.2, 0.35,
         font_size=10, color=ACCENT_RED, bold=True)

# No branch - continues
text_box(slide, "NO ↓", 5.9, 2.85, 2.0, 0.32,
         font_size=10, color=ACCENT_GREEN, bold=True)

box(slide, 4.8, 3.25, 3.73, 0.55, MID_GREY)
text_box(slide, "Whitespace discount on?", 4.9, 3.3, 3.55, 0.4,
         font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Yes - apply discount
box(slide, 8.8, 3.25, 2.8, 0.55, CARD_BG)
text_box(slide, "effectiveLOC =\nadditions × 0.95", 8.8, 3.25, 2.8, 0.55,
         font_size=10, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
text_box(slide, "YES →", 8.4, 3.35, 0.8, 0.3,
         font_size=10, color=ACCENT_ORANGE, bold=True)

# No - use raw
box(slide, 1.5, 3.25, 2.8, 0.55, CARD_BG)
text_box(slide, "effectiveLOC =\nadditions (raw)", 1.5, 3.25, 2.8, 0.55,
         font_size=10, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
text_box(slide, "← NO", 4.5, 3.35, 0.8, 0.3,
         font_size=10, color=ACCENT_TEAL, bold=True)

# Converge
text_box(slide, "↓", 6.23, 3.85, 0.4, 0.35, font_size=18, color=WHITE)
box(slide, 4.8, 4.25, 3.73, 0.55, ACCENT_TEAL)
text_box(slide, "Aggregate to DeveloperStats", 4.8, 4.25, 3.73, 0.55,
         font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

text_box(slide, "↓", 6.23, 4.85, 0.4, 0.35, font_size=18, color=WHITE)
box(slide, 4.8, 5.25, 3.73, 0.55, ACCENT_VIOLET)
text_box(slide, "LOC-burst pass: flag outliers", 4.8, 5.25, 3.73, 0.55,
         font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

text_box(slide, "↓", 6.23, 5.85, 0.4, 0.35, font_size=18, color=WHITE)
box(slide, 4.8, 6.25, 3.73, 0.5, ACCENT_ORANGE)
text_box(slide, "totalLOC / manualLOC / aiLOC", 4.8, 6.25, 3.73, 0.5,
         font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Side notes
notes = [
    ("Why additions only?",
     "Additions represent new code written. Deletions represent removed code "
     "and would undercount total effort.",
     "1.0", "4.5", ACCENT_BLUE),
    ("Why 5% discount?",
     "Conservative approximation. Exact blank-line filtering needs full diff "
     "fetch (100+ API calls for large repos).",
     "9.3", "4.5", ACCENT_ORANGE),
]
for title, desc, lx, ty, color in notes:
    lx = float(lx)
    ty = float(ty)
    box(slide, lx, ty, 3.5, 1.3, RGBColor(0x14, 0x20, 0x35))
    text_box(slide, title, lx+0.15, ty+0.08, 3.2, 0.3,
             font_size=10, bold=True, color=color)
    text_box(slide, desc, lx+0.15, ty+0.42, 3.2, 0.82,
             font_size=9, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DEPLOYMENT & OPERATIONS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Deployment & Operations",
             "Static SPA — deploy in minutes, zero ongoing infrastructure cost",
             ACCENT_GREEN)

# Build steps
text_box(slide, "Build & Deploy", 0.5, 1.38, 4.0, 0.4,
         font_size=14, bold=True, color=WHITE)

steps = [
    ("1", "npm install", "Install dependencies\n(React, Recharts, Tailwind)", ACCENT_BLUE),
    ("2", "npm run build", "Vite compiles to dist/\nMinified + tree-shaken", ACCENT_TEAL),
    ("3", "Upload dist/", "Copy to web host\nor CDN bucket", ACCENT_VIOLET),
    ("4", "Open in browser", "No config files\nNo environment vars", ACCENT_GREEN),
]
for i, (num, cmd, desc, color) in enumerate(steps):
    lx = 0.5 + i * 3.1
    box(slide, lx, 1.85, 2.8, 1.7, CARD_BG)
    accent_bar(slide, lx, 1.85, 2.8, 0.06, color)
    pill(slide, num, lx+0.15, 2.0, 0.45, 0.38, color, DARK_BG, 12)
    text_box(slide, cmd, lx+0.15, 2.52, 2.5, 0.38,
             font_size=12, bold=True, color=color)
    text_box(slide, desc, lx+0.15, 2.93, 2.5, 0.55,
             font_size=10, color=LIGHT_GREY)

# Deployment targets
text_box(slide, "Supported Deployment Targets", 0.5, 3.78, 5.0, 0.4,
         font_size=14, bold=True, color=WHITE)

targets = [
    ("☁️  AWS S3 + CloudFront", "Zero-config static hosting,\nglobal CDN, HTTPS included"),
    ("▲  Vercel / Netlify",     "Push to branch = auto-deploy,\npreviews for every PR"),
    ("🦑  GitHub / GitLab Pages","Free hosting, CI/CD built-in,\nideal for internal tools"),
    ("🐳  Docker + nginx",       "Self-hosted on-prem,\nbehind VPN or internal DNS"),
    ("📁  Local Filesystem",     "Open dist/index.html\ndirectly in browser (read-only)"),
    ("🖧   Internal Web Server",  "nginx / Apache serve dist/\nwith optional auth proxy"),
]
for i, (name, desc) in enumerate(targets):
    col = i % 3
    row = i // 3
    lx = 0.5 + col * 4.15
    ty = 4.25 + row * 1.55
    box(slide, lx, ty, 3.85, 1.35, CARD_BG)
    text_box(slide, name, lx+0.15, ty+0.1, 3.5, 0.38,
             font_size=12, bold=True, color=ACCENT_BLUE)
    text_box(slide, desc, lx+0.15, ty+0.52, 3.5, 0.75,
             font_size=10, color=LIGHT_GREY)

# Non-requirements
box(slide, 12.6, 1.35, 0.5, 5.9, CARD_BG)   # right margin

text_box(slide,
    "NOT REQUIRED:  database · backend server · cloud account · domain name · DevOps engineer",
    0.5, 7.15, 12.33, 0.3,
    font_size=10, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — USE CASES FOR LEADERSHIP
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Leadership Use Cases",
             "Practical questions this dashboard answers for engineering and product leaders",
             ACCENT_ORANGE)

use_cases = [
    ("📋", "Quarterly Business Reviews",
     "Show AI adoption rate trending over Q1→Q2→Q3. Correlate AI LOC growth "
     "with sprint velocity improvements. Demonstrate tool ROI to budget holders.",
     ACCENT_BLUE, [
         "Filter by 3-month ranges",
         "Compare Manual vs AI LOC quarter-on-quarter",
         "Export to slide by screenshot",
     ]),
    ("👥", "Team Health & Coaching",
     "Identify developers with 0% AI adoption (may need tooling training) and those "
     "with 80%+ AI LOC (verify quality standards are met and no over-reliance).",
     ACCENT_TEAL, [
         "Per-developer AI % column in table",
         "Commit drill-down for context",
         "Focus on LOC burst anomalies",
     ]),
    ("🔒", "Compliance & Audit",
     "For regulated industries, document which code was AI-assisted for IP, "
     "licensing, or compliance reviews. Every AI commit has a traceable reason tag.",
     ACCENT_VIOLET, [
         "aiReason field: pattern / loc-burst / committer-mismatch",
         "Drill to individual commit SHAs",
         "Date-range filter for audit period",
     ]),
    ("⚡", "Sprint Retrospectives",
     "In 30-second sprint reviews: pull last 14-day data, see which engineers "
     "used AI, and discuss if AI acceleration matched sprint commitments.",
     ACCENT_ORANGE, [
         "14-day preset in date picker",
         "Developer table sorted by AI LOC",
         "Live fetch — always fresh data",
     ]),
    ("🎯", "AI Tool ROI Measurement",
     "Compare AI LOC percentage before and after rolling out a new tool (e.g., Copilot). "
     "Validate adoption is actually happening, not just license spend.",
     ACCENT_GREEN, [
         "Pattern filter catches tool-specific tags",
         "Date range before vs after rollout",
         "Per-developer granularity",
     ]),
    ("📊", "Engineering KPI Dashboards",
     "Embed the AI% metric into your existing OKR tracking. "
     "Goal: '>30% of new code AI-assisted by Q4' is now measurable and auditable.",
     ACCENT_RED, [
         "5 KPI cards on dashboard home",
         "aiPercentage = aiLOC / totalLOC × 100",
         "Historical snapshots via screenshots",
     ]),
]

for i, (icon, title, desc, color, bullets) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    lx = 0.38 + col * 4.32
    ty = 1.35 + row * 2.88
    card = box(slide, lx, ty, 4.05, 2.65, CARD_BG)
    accent_bar(slide, lx, ty, 4.05, 0.06, color)
    text_box(slide, icon, lx+0.15, ty+0.12, 0.5, 0.45, font_size=18)
    text_box(slide, title, lx+0.65, ty+0.15, 3.2, 0.38,
             font_size=12, bold=True, color=color)
    text_box(slide, desc, lx+0.15, ty+0.62, 3.75, 0.9,
             font_size=10, color=LIGHT_GREY)
    for j, b in enumerate(bullets):
        text_box(slide, f"› {b}", lx+0.15, ty+1.58+j*0.34, 3.75, 0.32,
                 font_size=9, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — KNOWN LIMITATIONS & ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide)
slide_header(slide, "Limitations & Potential Roadmap",
             "Being transparent about current tradeoffs and future opportunities",
             ACCENT_ORANGE)

# Current limitations
text_box(slide, "Current Limitations", 0.5, 1.38, 5.8, 0.4,
         font_size=14, bold=True, color=ACCENT_ORANGE)

limitations = [
    ("5% whitespace approximation",
     "Blank-line accuracy requires fetching full diffs — too expensive at scale. "
     "Current 5% discount is a conservative engineering tradeoff."),
    ("Sequential project fetching",
     "Projects processed one at a time to respect GitLab's 10 req/s rate limit for "
     "personal tokens. OAuth or service accounts could enable parallel fetching."),
    ("No historical snapshots",
     "Each fetch is independent. Trending metrics (week-over-week) require "
     "manually capturing screenshots or adding a lightweight persistence layer."),
    ("Client-side author filtering",
     "GitLab API only supports single-author filter, so all commits are fetched "
     "then filtered in-browser. Large groups may load slowly."),
    ("No co-author aggregation",
     "Human co-authors from git trailers are not attributed separately. "
     "Only bot co-authors are used as an AI signal."),
]
for i, (title, desc) in enumerate(limitations):
    ty = 1.87 + i * 0.9
    box(slide, 0.5, ty, 6.0, 0.78, CARD_BG)
    pill(slide, "!", 0.65, ty+0.18, 0.32, 0.32, ACCENT_ORANGE, DARK_BG, 11)
    text_box(slide, title, 1.12, ty+0.08, 5.2, 0.3,
             font_size=11, bold=True, color=ACCENT_ORANGE)
    text_box(slide, desc, 1.12, ty+0.4, 5.2, 0.35,
             font_size=9, color=LIGHT_GREY)

# Roadmap
text_box(slide, "Potential Roadmap", 7.0, 1.38, 6.0, 0.4,
         font_size=14, bold=True, color=ACCENT_GREEN)

roadmap = [
    ("Near Term", ACCENT_BLUE, [
        "GitHub / Bitbucket / Azure DevOps support",
        "CSV / Excel export of developer stats",
        "Per-project breakdown (not just per-developer)",
        "Configurable LOC burst multiplier in UI",
    ]),
    ("Medium Term", ACCENT_VIOLET, [
        "LocalStorage snapshot history (opt-in)",
        "Trend charts: AI% over rolling weeks",
        "Team / squad grouping config",
        "Shareable URL with encoded config",
    ]),
    ("Longer Term", ACCENT_TEAL, [
        "REST API backend for persistent data warehouse",
        "Slack / Teams weekly digest notifications",
        "AI tool attribution analytics (Copilot vs Claude vs Cursor)",
        "Diff-level whitespace analysis",
    ]),
]
for i, (phase, color, items) in enumerate(roadmap):
    ty = 1.87 + i * 1.88
    box(slide, 7.0, ty, 6.0, 1.72, CARD_BG)
    accent_bar(slide, 7.0, ty, 6.0, 0.06, color)
    text_box(slide, phase, 7.15, ty+0.12, 2.0, 0.35,
             font_size=12, bold=True, color=color)
    for j, item in enumerate(items):
        text_box(slide, f"  › {item}", 7.15, ty+0.52+j*0.3, 5.65, 0.3,
                 font_size=10, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SUMMARY & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide, DARKER_CARD)

# Background accent
corner = slide.shapes.add_shape(1,
    Inches(8.5), Inches(3.5), Inches(6.0), Inches(6.0))
corner.fill.solid()
corner.fill.fore_color.rgb = RGBColor(0x1A, 0x26, 0x3D)
corner.line.fill.background()

accent_bar(slide, 0, 0, 0.12, 7.5, ACCENT_TEAL)
accent_bar(slide, 0, 0, 13.33, 0.07, ACCENT_TEAL)

text_box(slide, "Summary", 0.5, 0.2, 8.0, 0.55,
         font_size=28, bold=True, color=WHITE)

# 5 key takeaways
takeaways = [
    ("1", "Real-time AI adoption visibility",
     "No surveys. No manual auditing. Instant data from your GitLab.",
     ACCENT_TEAL),
    ("2", "3-signal detection accuracy",
     "Commit patterns, tool committer mismatch, and LOC-burst heuristic work together.",
     ACCENT_VIOLET),
    ("3", "Zero infrastructure footprint",
     "Static SPA. No server. No database. Deploy in < 5 minutes.",
     ACCENT_BLUE),
    ("4", "Security-first token handling",
     "Token lives only in RAM. Cleared on refresh. Never persisted or logged.",
     ACCENT_GREEN),
    ("5", "Leadership-ready out of the box",
     "KPI cards, drill-down tables, and charts designed for executive briefings.",
     ACCENT_ORANGE),
]
for i, (num, title, desc, color) in enumerate(takeaways):
    ty = 1.0 + i * 1.12
    box(slide, 0.5, ty, 8.0, 1.0, CARD_BG)
    accent_bar(slide, 0.5, ty, 8.0, 0.05, color)
    pill(slide, num, 0.65, ty+0.27, 0.4, 0.38, color, DARK_BG, 13)
    text_box(slide, title, 1.2, ty+0.1, 7.1, 0.38,
             font_size=13, bold=True, color=color)
    text_box(slide, desc, 1.2, ty+0.5, 7.1, 0.38,
             font_size=11, color=LIGHT_GREY)

# Next steps CTA
box(slide, 8.8, 1.0, 4.3, 4.3, CARD_BG)
accent_bar(slide, 8.8, 1.0, 4.3, 0.07, ACCENT_ORANGE)
text_box(slide, "Next Steps", 9.0, 1.15, 3.9, 0.45,
         font_size=15, bold=True, color=ACCENT_ORANGE)
next_steps = [
    "1.  Get a GitLab PAT (read_api scope)",
    "2.  Note your Group or Project ID",
    "3.  Clone or download the repo",
    "4.  Run  npm install && npm run dev",
    "5.  Open  http://localhost:3000",
    "6.  Enter URL, token, and IDs",
    "7.  Click Fetch Data",
    "8.  Share findings with your team!",
]
multi_para(slide, next_steps, 9.0, 1.7, 3.9, 3.5,
           font_size=11, color=LIGHT_GREY, spacing_after=8)

# Bottom bar
box(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x0A, 0x10, 0x20))
text_box(slide,
    "GitAnalytics AI Dashboard  ·  April 2026  ·  Built with React 18 + TypeScript + Recharts  ·  Read-only · Zero backend · No data persisted",
    0.0, 7.08, 13.33, 0.35,
    font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
out_path = "/Users/788763/gitanalytics-dashboard/GitAnalytics_AI_Dashboard_Leadership.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
